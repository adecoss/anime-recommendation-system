from __future__ import annotations

from pathlib import Path
import textwrap

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
REPORTS = ROOT / "reports"
PLOTS = ROOT / "artifacts" / "plots"
RECO = ROOT / "artifacts" / "recommendation"
OUTPUT = REPORTS / "week10_recommendation_presentation.pptx"


TITLE = RGBColor(32, 47, 67)
BLUE = RGBColor(76, 120, 168)
ORANGE = RGBColor(242, 142, 43)
GREEN = RGBColor(89, 161, 79)
PURPLE = RGBColor(176, 122, 161)
GRAY = RGBColor(90, 90, 90)
LIGHT = RGBColor(245, 247, 250)


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.35), Inches(0.25), Inches(12.6), Inches(0.45))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = TITLE
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.36), Inches(0.72), Inches(12.4), Inches(0.3))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(10.5)
        sp.font.color.rgb = GRAY


def add_bullets(slide, items: list[str], left: float, top: float, width: float, height: float, font_size: int = 15) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = TITLE
        p.space_after = Pt(6)


def add_card(slide, title: str, body: str, left: float, top: float, width: float, height: float, fill=LIGHT) -> None:
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(210, 215, 222)
    tf = shape.text_frame
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.10)
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = TITLE
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(10.5)
    p2.font.color.rgb = GRAY


def add_picture_if_exists(slide, path: Path, left: float, top: float, width: float | None = None, height: float | None = None) -> None:
    if path.exists():
        kwargs = {}
        if width is not None:
            kwargs["width"] = Inches(width)
        if height is not None:
            kwargs["height"] = Inches(height)
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), **kwargs)


def add_table(slide, df: pd.DataFrame, left: float, top: float, width: float, height: float, font_size: int = 8) -> None:
    rows = len(df) + 1
    cols = len(df.columns)
    table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)).table
    for col_idx, col in enumerate(df.columns):
        table.cell(0, col_idx).text = str(col)
        table.cell(0, col_idx).fill.solid()
        table.cell(0, col_idx).fill.fore_color.rgb = BLUE
        for paragraph in table.cell(0, col_idx).text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
    for row_idx, (_, row) in enumerate(df.iterrows(), start=1):
        for col_idx, value in enumerate(row):
            text = str(value)
            if len(text) > 95:
                text = text[:92] + "..."
            table.cell(row_idx, col_idx).text = text
            for paragraph in table.cell(row_idx, col_idx).text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.color.rgb = TITLE


def short(text: str, width: int = 54) -> str:
    return "\n".join(textwrap.wrap(str(text), width=width))


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    metrics = pd.read_csv(RECO / "week10_evaluation_metrics.csv")
    level_metrics = pd.read_csv(RECO / "week10_metrics_by_user_level.csv")
    beginner = pd.read_csv(RECO / "week10_beginner_entrypoint_candidates.csv")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Anime Discovery, Recommendation, and Graph Intelligence", "Week 10 deliverable: recommendation, ranking, and evaluation")
    add_bullets(
        slide,
        [
            "Goal: rank anime that a user is likely to enjoy next.",
            "Built on a reproducible catalog, interaction layer, graph relations, and feature matrices.",
            "Week 10 adds offline recommendation evaluation and a product-level recommender design.",
        ],
        0.7,
        1.35,
        6.0,
        2.1,
        18,
    )
    add_card(slide, "Main decision", "Given liked/completed anime, choose the next titles to recommend.", 7.2, 1.25, 4.9, 1.05)
    add_card(slide, "Not claimed", "This is not chronological next-watch prediction because the ratings source has no timestamps.", 7.2, 2.55, 4.9, 1.05)
    add_card(slide, "Core result", "Personalized SVD and hybrid ranking both beat popularity-only discovery.", 7.2, 3.85, 4.9, 1.05)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "From Dataset to Recommender", "The Week 10 model uses all previous project layers")
    add_card(slide, "Catalog layer", "MAL/Jikan metadata enriched with AniDB and Shoko XML cache.", 0.55, 1.2, 3.0, 1.3)
    add_card(slide, "Feature layer", "Numeric, categorical, text, tag, runtime, and graph-derived fields.", 3.85, 1.2, 3.0, 1.3)
    add_card(slide, "Interaction layer", "Anonymized user-anime ratings filtered into positive interactions.", 7.15, 1.2, 3.0, 1.3)
    add_card(slide, "Graph layer", "Recommendation edges plus typed relation edges for franchise navigation.", 10.45, 1.2, 2.35, 1.3)
    add_picture_if_exists(slide, PLOTS / "eda" / "bakemonogatari_relation_tree.png", 1.0, 3.0, width=11.2)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Representation and Segmentation Recap", "Week 5 and Week 7 are inputs, not isolated exercises")
    add_picture_if_exists(slide, PLOTS / "week5" / "svd_sparse_catalog_energy.png", 0.55, 1.15, width=5.8)
    add_picture_if_exists(slide, PLOTS / "week7" / "cluster_genre_profile_heatmap.png", 6.85, 1.15, width=5.85)
    add_bullets(
        slide,
        [
            "SVD gives a compact catalog representation for similarity and clustering.",
            "Clusters are broad discovery segments, not hard genre truth.",
            "Density methods showed local pockets but rejected too many titles as noise.",
        ],
        0.8,
        6.25,
        11.8,
        0.7,
        12,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "User Levels and Controls", "The final recommender should adapt to user maturity")
    add_card(slide, "Absolute Beginner", "Onboarding, genre choice, recognizable short entry points.", 0.55, 1.15, 2.95, 1.25, RGBColor(236, 244, 255))
    add_card(slide, "Amateur", "Similar anime, franchise relations, high-rated anchors.", 3.75, 1.15, 2.95, 1.25, RGBColor(238, 248, 239))
    add_card(slide, "Good", "Collaborative ranking plus seasonal/current exploration.", 6.95, 1.15, 2.95, 1.25, RGBColor(255, 246, 232))
    add_card(slide, "Pro", "Long-tail, novelty, niche clusters, graph-aware exploration.", 10.15, 1.15, 2.65, 1.25, RGBColor(248, 239, 248))
    add_bullets(
        slide,
        [
            "User controls: genre, demographic/content rating, explicit toggle, score floor, episode length, year.",
            "Controls constrain the candidate set; ranking still decides the order.",
            "Relations should offer franchise continuation after a main entry is selected, not blindly recommend season 3 first.",
        ],
        1.0,
        3.15,
        11.3,
        1.4,
        17,
    )
    add_picture_if_exists(slide, PLOTS / "week10" / "week10_user_level_distribution.png", 2.5, 5.0, width=8.2)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Evaluation Design", "Ranking recovery test with a serious popularity baseline")
    eval_cards = [
        ("Positive", "rating >= 7"),
        ("Candidate pool", "8,181 anime with enough positive ratings"),
        ("Evaluation users", "15,000 users"),
        ("Candidate list", "1 held-out liked item + 100 sampled negatives"),
    ]
    for idx, (t, b) in enumerate(eval_cards):
        add_card(slide, t, b, 0.65 + idx * 3.15, 1.15, 2.75, 1.2)
    add_bullets(
        slide,
        [
            "Held-out user-item pairs are removed from the training matrix.",
            "No timestamp claims: this evaluates preference recovery, not exact next-watch sequence.",
            "Popularity is not a strawman in anime; it reflects real discourse and entry-point behavior.",
        ],
        0.9,
        3.0,
        11.8,
        1.3,
        17,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Overall Recommendation Results", "Personalized models beat popularity-only discovery")
    table_df = metrics.copy()
    table_df["hit_rate_at_10"] = table_df["hit_rate_at_10"].map(lambda x: f"{x:.3f}")
    table_df["ndcg_at_10"] = table_df["ndcg_at_10"].map(lambda x: f"{x:.3f}")
    table_df["mean_reciprocal_rank"] = table_df["mean_reciprocal_rank"].map(lambda x: f"{x:.3f}")
    table_df = table_df[["method", "hit_rate_at_10", "ndcg_at_10", "mean_reciprocal_rank", "median_rank"]]
    add_table(slide, table_df, 0.7, 1.1, 5.95, 1.5, font_size=8)
    add_picture_if_exists(slide, PLOTS / "week10" / "week10_metric_comparison.png", 6.95, 1.0, width=5.8)
    add_picture_if_exists(slide, PLOTS / "week10" / "week10_median_rank.png", 3.2, 4.0, width=6.8)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Results by User Level", "The same model has different meaning for different users")
    hybrid = level_metrics[level_metrics["method"].eq("hybrid_svd_popularity")].copy()
    hybrid["hit_rate_at_10"] = hybrid["hit_rate_at_10"].map(lambda x: f"{x:.3f}")
    hybrid = hybrid[["user_level", "evaluated_users", "median_profile_size", "hit_rate_at_10", "median_rank"]]
    add_table(slide, hybrid, 0.65, 1.0, 6.0, 1.8, font_size=8)
    add_picture_if_exists(slide, PLOTS / "week10" / "week10_hit_at_10_by_user_level.png", 7.1, 0.95, width=5.55)
    add_bullets(
        slide,
        [
            "Amateur users perform best: enough signal, but not too much catalog saturation.",
            "Good users are still strong but more diverse.",
            "Pro users are rare and harder; this segment needs novelty and long-tail logic.",
        ],
        1.0,
        4.55,
        11.4,
        1.2,
        16,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Absolute Beginner Entry Points", "Cold-start pool before collaborative ranking is reliable")
    top = beginner.head(8)[["title", "type", "score", "episodes", "genres"]].copy()
    top["genres"] = top["genres"].apply(lambda x: short(x, 36))
    add_table(slide, top, 0.6, 1.0, 7.1, 4.8, font_size=7)
    add_bullets(
        slide,
        [
            "Filters out explicit titles and obvious continuation entries.",
            "Prefers high score, high engagement, and shorter or medium-length commitments.",
            "Used for onboarding after genre/content choices.",
        ],
        8.0,
        1.15,
        4.65,
        2.6,
        16,
    )
    add_card(slide, "Example behavior", "If a beginner likes Frieren or Steins;Gate, the system can then switch into similar-anime and graph-aware recommendations.", 8.0, 4.15, 4.65, 1.25)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Recommendation Architecture", "Use the right signal at the right moment")
    add_bullets(
        slide,
        [
            "1. Apply user controls: genre, rating, explicit toggle, score, year, length.",
            "2. Select maturity mode: beginner, amateur, good, or pro.",
            "3. Generate candidates: popularity, SVD, content similarity, relation graph.",
            "4. Rank with hybrid score.",
            "5. Use relations for franchise navigation after an entry point is chosen.",
        ],
        0.9,
        1.1,
        6.0,
        3.2,
        17,
    )
    add_picture_if_exists(slide, PLOTS / "eda" / "recommendation_network_top_popular.png", 7.0, 1.0, width=5.7)
    add_card(slide, "Design rule", "General discovery should favor standalone or first-entry titles. Sequels and side stories belong in a contextual continuation rail.", 1.2, 5.2, 10.7, 0.9, RGBColor(255, 246, 232))

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Limitations and Defense Points", "What the current result supports, and what it does not")
    add_bullets(
        slide,
        [
            "Supports: collaborative ranking improves over popularity-only recommendation.",
            "Supports: user maturity matters; the Pro segment needs a different strategy.",
            "Supports: relation edges are necessary for franchise-aware recommendation.",
            "Does not prove: chronological next-watch prediction.",
            "Does not solve yet: fresh seasonal recency, full watch-order data, or perfect MAL/AniDB split handling.",
        ],
        0.9,
        1.15,
        11.7,
        3.2,
        18,
    )
    add_card(slide, "Next technical focus", "Week 12 can move graph analytics from diagnostics into graph-aware reranking: centrality, relation paths, franchise entry points, and novelty constraints.", 1.0, 5.35, 11.2, 1.0, RGBColor(236, 244, 255))

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"wrote {OUTPUT}")


if __name__ == "__main__":
    main()
